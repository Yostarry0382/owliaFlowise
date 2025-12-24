"""
PowerPoint Generation Service for Flowise Integration
python-pptx を使用してテンプレートベースのスライド生成を行う
"""

import os
import json
import uuid
import shutil
from datetime import datetime
from typing import Optional, List, Dict, Any
from pathlib import Path

from fastapi import FastAPI, HTTPException, UploadFile, File, Form
from fastapi.responses import FileResponse
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel, Field
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RgbColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE_TYPE


# ===== 設定 =====
BASE_DIR = Path(__file__).parent
TEMPLATES_DIR = BASE_DIR / "templates"
OUTPUT_DIR = BASE_DIR / "output"
TEMP_DIR = BASE_DIR / "temp"

# ディレクトリ作成
TEMPLATES_DIR.mkdir(exist_ok=True)
OUTPUT_DIR.mkdir(exist_ok=True)
TEMP_DIR.mkdir(exist_ok=True)


# ===== Pydantic Models =====

class SlideContent(BaseModel):
    """個別スライドのコンテンツ定義"""
    layout_index: int = Field(default=0, description="使用するレイアウトのインデックス")
    title: Optional[str] = Field(default=None, description="スライドタイトル")
    subtitle: Optional[str] = Field(default=None, description="サブタイトル")
    body: Optional[str] = Field(default=None, description="本文テキスト")
    bullets: Optional[List[str]] = Field(default=None, description="箇条書きリスト")
    notes: Optional[str] = Field(default=None, description="発表者ノート")
    image_path: Optional[str] = Field(default=None, description="挿入する画像のパス")
    placeholders: Optional[Dict[int, str]] = Field(
        default=None,
        description="プレースホルダーインデックスとテキストのマッピング"
    )


class PresentationRequest(BaseModel):
    """プレゼンテーション生成リクエスト"""
    template_id: Optional[str] = Field(default=None, description="使用するテンプレートID")
    slides: List[SlideContent] = Field(..., description="スライドコンテンツのリスト")
    output_filename: Optional[str] = Field(default=None, description="出力ファイル名")
    metadata: Optional[Dict[str, str]] = Field(default=None, description="メタデータ（作成者など）")


class TemplateInfo(BaseModel):
    """テンプレート情報"""
    id: str
    name: str
    description: Optional[str]
    layouts: List[Dict[str, Any]]
    created_at: str
    file_path: str


class SlideAnalysis(BaseModel):
    """スライド解析結果"""
    slide_index: int
    layout_name: str
    placeholders: List[Dict[str, Any]]
    shapes: List[Dict[str, Any]]


class TemplateAnalysis(BaseModel):
    """テンプレート解析結果"""
    template_id: str
    slide_masters: List[Dict[str, Any]]
    layouts: List[Dict[str, Any]]
    slides: List[SlideAnalysis]


# ===== FastAPI App =====

app = FastAPI(
    title="PPTX Generator Service",
    description="PowerPoint生成サービス - Flowise連携用",
    version="1.0.0"
)

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)


# ===== ヘルパー関数 =====

def get_template_path(template_id: str) -> Path:
    """テンプレートファイルのパスを取得"""
    template_path = TEMPLATES_DIR / f"{template_id}.pptx"
    if not template_path.exists():
        raise HTTPException(status_code=404, detail=f"Template not found: {template_id}")
    return template_path


def analyze_placeholder(placeholder) -> Dict[str, Any]:
    """プレースホルダーを解析"""
    return {
        "idx": placeholder.placeholder_format.idx,
        "type": str(placeholder.placeholder_format.type),
        "name": placeholder.name,
        "width": placeholder.width.inches if placeholder.width else None,
        "height": placeholder.height.inches if placeholder.height else None,
        "left": placeholder.left.inches if placeholder.left else None,
        "top": placeholder.top.inches if placeholder.top else None,
    }


def analyze_shape(shape) -> Dict[str, Any]:
    """シェイプを解析"""
    info = {
        "name": shape.name,
        "shape_type": str(shape.shape_type),
        "has_text_frame": shape.has_text_frame,
    }
    if shape.has_text_frame:
        info["text"] = shape.text_frame.text[:100] if shape.text_frame.text else ""
    return info


def set_text_in_placeholder(placeholder, text: str, font_size: Optional[int] = None):
    """プレースホルダーにテキストを設定"""
    if placeholder.has_text_frame:
        tf = placeholder.text_frame
        # 既存のテキストをクリア
        for paragraph in tf.paragraphs:
            paragraph.clear()
        # 新しいテキストを設定
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = text
        if font_size:
            run.font.size = Pt(font_size)


def add_bullets_to_placeholder(placeholder, bullets: List[str]):
    """プレースホルダーに箇条書きを追加"""
    if placeholder.has_text_frame:
        tf = placeholder.text_frame
        # 最初の段落を使用
        for i, bullet in enumerate(bullets):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = bullet
            p.level = 0


# ===== API Endpoints =====

@app.get("/")
async def root():
    """ヘルスチェック"""
    return {
        "status": "healthy",
        "service": "PPTX Generator",
        "version": "1.0.0"
    }


@app.get("/templates", response_model=List[TemplateInfo])
async def list_templates():
    """登録済みテンプレート一覧を取得"""
    templates = []
    for pptx_file in TEMPLATES_DIR.glob("*.pptx"):
        template_id = pptx_file.stem
        try:
            prs = Presentation(str(pptx_file))
            layouts = []
            for i, layout in enumerate(prs.slide_layouts):
                layouts.append({
                    "index": i,
                    "name": layout.name
                })

            templates.append(TemplateInfo(
                id=template_id,
                name=template_id,
                description=f"Template: {template_id}",
                layouts=layouts,
                created_at=datetime.fromtimestamp(pptx_file.stat().st_mtime).isoformat(),
                file_path=str(pptx_file)
            ))
        except Exception as e:
            print(f"Error loading template {template_id}: {e}")

    return templates


@app.post("/templates/upload")
async def upload_template(
    file: UploadFile = File(...),
    template_id: Optional[str] = Form(None),
    description: Optional[str] = Form(None)
):
    """テンプレートをアップロード"""
    if not file.filename.endswith('.pptx'):
        raise HTTPException(status_code=400, detail="Only .pptx files are supported")

    # テンプレートIDを生成または使用
    tid = template_id or Path(file.filename).stem
    template_path = TEMPLATES_DIR / f"{tid}.pptx"

    # ファイルを保存
    with open(template_path, "wb") as f:
        content = await file.read()
        f.write(content)

    # テンプレート情報を返す
    prs = Presentation(str(template_path))
    layouts = []
    for i, layout in enumerate(prs.slide_layouts):
        layouts.append({
            "index": i,
            "name": layout.name
        })

    return {
        "message": "Template uploaded successfully",
        "template_id": tid,
        "layouts": layouts
    }


@app.get("/templates/{template_id}/analyze", response_model=TemplateAnalysis)
async def analyze_template(template_id: str):
    """テンプレートの構造を解析"""
    template_path = get_template_path(template_id)
    prs = Presentation(str(template_path))

    # スライドマスター情報
    slide_masters = []
    for i, master in enumerate(prs.slide_masters):
        slide_masters.append({
            "index": i,
            "name": master.name if hasattr(master, 'name') else f"Master {i}"
        })

    # レイアウト情報
    layouts = []
    for i, layout in enumerate(prs.slide_layouts):
        placeholders = []
        for ph in layout.placeholders:
            placeholders.append(analyze_placeholder(ph))
        layouts.append({
            "index": i,
            "name": layout.name,
            "placeholders": placeholders
        })

    # 既存スライド情報
    slides = []
    for i, slide in enumerate(prs.slides):
        placeholders = []
        shapes = []
        for shape in slide.shapes:
            if shape.is_placeholder:
                placeholders.append(analyze_placeholder(shape))
            shapes.append(analyze_shape(shape))

        slides.append(SlideAnalysis(
            slide_index=i,
            layout_name=slide.slide_layout.name,
            placeholders=placeholders,
            shapes=shapes
        ))

    return TemplateAnalysis(
        template_id=template_id,
        slide_masters=slide_masters,
        layouts=layouts,
        slides=slides
    )


@app.post("/generate")
async def generate_presentation(request: PresentationRequest):
    """プレゼンテーションを生成"""

    # テンプレートを読み込むか新規作成
    if request.template_id:
        template_path = get_template_path(request.template_id)
        prs = Presentation(str(template_path))
        # テンプレートの既存スライドを削除（レイアウトのみ使用）
        while len(prs.slides) > 0:
            rId = prs.slides._sldIdLst[0].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[0]
    else:
        prs = Presentation()

    # メタデータを設定
    if request.metadata:
        core_props = prs.core_properties
        if "author" in request.metadata:
            core_props.author = request.metadata["author"]
        if "title" in request.metadata:
            core_props.title = request.metadata["title"]
        if "subject" in request.metadata:
            core_props.subject = request.metadata["subject"]

    # スライドを追加
    for slide_content in request.slides:
        layout_index = slide_content.layout_index
        if layout_index >= len(prs.slide_layouts):
            layout_index = 0

        slide_layout = prs.slide_layouts[layout_index]
        slide = prs.slides.add_slide(slide_layout)

        # タイトルを設定
        if slide_content.title and slide.shapes.title:
            slide.shapes.title.text = slide_content.title

        # プレースホルダーにコンテンツを設定
        for shape in slide.shapes:
            if shape.is_placeholder:
                ph_idx = shape.placeholder_format.idx
                ph_type = str(shape.placeholder_format.type)

                # サブタイトル（通常idx=1）
                if slide_content.subtitle and ph_idx == 1:
                    set_text_in_placeholder(shape, slide_content.subtitle)

                # 本文テキスト
                elif slide_content.body and "BODY" in ph_type:
                    set_text_in_placeholder(shape, slide_content.body)

                # 箇条書き
                elif slide_content.bullets and "BODY" in ph_type:
                    add_bullets_to_placeholder(shape, slide_content.bullets)

                # カスタムプレースホルダーマッピング
                elif slide_content.placeholders and ph_idx in slide_content.placeholders:
                    set_text_in_placeholder(shape, slide_content.placeholders[ph_idx])

        # 発表者ノート
        if slide_content.notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = slide_content.notes

    # ファイルを保存
    output_filename = request.output_filename or f"presentation_{uuid.uuid4().hex[:8]}.pptx"
    if not output_filename.endswith('.pptx'):
        output_filename += '.pptx'
    output_path = OUTPUT_DIR / output_filename
    prs.save(str(output_path))

    return {
        "success": True,
        "message": "Presentation generated successfully",
        "filename": output_filename,
        "download_url": f"/download/{output_filename}",
        "slide_count": len(prs.slides)
    }


@app.post("/generate/from-json")
async def generate_from_json(
    json_content: str = Form(...),
    template_id: Optional[str] = Form(None)
):
    """JSON文字列からプレゼンテーションを生成（Flowise Tool用）"""
    try:
        data = json.loads(json_content)
        request = PresentationRequest(
            template_id=template_id,
            slides=[SlideContent(**slide) for slide in data.get("slides", [])],
            output_filename=data.get("output_filename"),
            metadata=data.get("metadata")
        )
        return await generate_presentation(request)
    except json.JSONDecodeError as e:
        raise HTTPException(status_code=400, detail=f"Invalid JSON: {e}")
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.get("/download/{filename}")
async def download_file(filename: str):
    """生成したファイルをダウンロード"""
    file_path = OUTPUT_DIR / filename
    if not file_path.exists():
        raise HTTPException(status_code=404, detail="File not found")

    return FileResponse(
        path=str(file_path),
        filename=filename,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )


@app.delete("/files/{filename}")
async def delete_file(filename: str):
    """生成したファイルを削除"""
    file_path = OUTPUT_DIR / filename
    if file_path.exists():
        file_path.unlink()
        return {"message": f"File {filename} deleted"}
    raise HTTPException(status_code=404, detail="File not found")


@app.post("/templates/{template_id}/fill")
async def fill_template(template_id: str, content: Dict[str, Any]):
    """
    テンプレートのスライドを維持しながらコンテンツを埋める
    既存スライドの構造を保持したまま、テキストのみ置換
    """
    template_path = get_template_path(template_id)
    prs = Presentation(str(template_path))

    slides_content = content.get("slides", [])

    for i, slide in enumerate(prs.slides):
        if i >= len(slides_content):
            break

        slide_data = slides_content[i]

        for shape in slide.shapes:
            if shape.is_placeholder:
                ph_idx = shape.placeholder_format.idx

                # タイトル（通常idx=0）
                if ph_idx == 0 and "title" in slide_data:
                    set_text_in_placeholder(shape, slide_data["title"])

                # サブタイトル/本文（通常idx=1以上）
                elif ph_idx == 1 and "subtitle" in slide_data:
                    set_text_in_placeholder(shape, slide_data["subtitle"])

                # カスタムマッピング
                elif "placeholders" in slide_data:
                    ph_key = str(ph_idx)
                    if ph_key in slide_data["placeholders"]:
                        set_text_in_placeholder(shape, slide_data["placeholders"][ph_key])

        # ノート
        if "notes" in slide_data:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = slide_data["notes"]

    # 保存
    output_filename = content.get("output_filename", f"filled_{template_id}_{uuid.uuid4().hex[:8]}.pptx")
    if not output_filename.endswith('.pptx'):
        output_filename += '.pptx'
    output_path = OUTPUT_DIR / output_filename
    prs.save(str(output_path))

    return {
        "success": True,
        "message": "Template filled successfully",
        "filename": output_filename,
        "download_url": f"/download/{output_filename}",
        "slide_count": len(prs.slides)
    }


if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8100)
